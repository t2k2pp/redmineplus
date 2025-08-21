from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import io
from datetime import datetime
from typing import Dict, List
import textwrap

class PowerPointGenerator:
    def __init__(self):
        self.prs = None
    
    def create_issue_report(self, issue_data: Dict) -> bytes:
        self.prs = Presentation()
        
        # 1ページ目のみ作成（画像に基づいてレイアウト）
        self._create_issue_report_page(issue_data)
        
        return self._save_to_bytes()
    
    def _create_issue_report_page(self, issue_data: Dict):
        # 空白レイアウトを使用
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # トラッカー・ID・ステータスヘッダー部分
        self._add_header_section(slide, issue_data)
        
        # タイトル（件名）
        self._add_title_section(slide, issue_data)
        
        # プロジェクト情報
        self._add_project_section(slide, issue_data)
        
        # 左側: 説明セクション
        self._add_description_section(slide, issue_data)
        
        # 右側: 詳細情報セクション
        self._add_details_section(slide, issue_data)
        
        # コメントセクション
        self._add_comments_section(slide, issue_data)
        
        # 更新履歴フッター
        self._add_footer_section(slide, issue_data)
    
    def _add_header_section(self, slide, issue_data: Dict):
        # Bugラベル（赤背景）
        tracker_name = issue_data.get('tracker', {}).get('name', 'チケット')
        bug_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(0.8), Inches(0.4)
        )
        bug_shape.fill.solid()
        bug_shape.fill.fore_color.rgb = RGBColor(220, 53, 69)  # 赤色
        bug_shape.line.color.rgb = RGBColor(220, 53, 69)
        bug_text = bug_shape.text_frame
        bug_text.text = tracker_name
        bug_text.paragraphs[0].font.size = Pt(12)
        bug_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        bug_text.paragraphs[0].font.bold = True
        
        # チケット番号
        ticket_id = f"#{issue_data.get('id', '')}"
        id_textbox = slide.shapes.add_textbox(
            Inches(1.5), Inches(0.3), Inches(1.5), Inches(0.4)
        )
        id_text = id_textbox.text_frame
        id_text.text = ticket_id
        id_text.paragraphs[0].font.size = Pt(16)
        id_text.paragraphs[0].font.color.rgb = RGBColor(108, 117, 125)  # グレー
        id_text.paragraphs[0].font.bold = True
        
        # ステータス（右側）
        status_name = issue_data.get('status', {}).get('name', '新規')
        status_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(0.3), Inches(1.0), Inches(0.4)
        )
        status_shape.fill.solid()
        status_shape.fill.fore_color.rgb = RGBColor(23, 162, 184)  # 青色
        status_shape.line.color.rgb = RGBColor(23, 162, 184)
        status_text = status_shape.text_frame
        status_text.text = status_name
        status_text.paragraphs[0].font.size = Pt(12)
        status_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        status_text.paragraphs[0].font.bold = True
    
    def _add_title_section(self, slide, issue_data: Dict):
        title_text = issue_data.get('subject', '')
        title_textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.6)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(33, 37, 41)
    
    def _add_project_section(self, slide, issue_data: Dict):
        project_name = issue_data.get('project', {}).get('name', '')
        project_textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(9.0), Inches(0.3)
        )
        project_frame = project_textbox.text_frame
        project_frame.text = f"プロジェクト: {project_name}"
        project_frame.paragraphs[0].font.size = Pt(12)
        project_frame.paragraphs[0].font.color.rgb = RGBColor(108, 117, 125)
    
    def _add_description_section(self, slide, issue_data: Dict):
        # 説明ヘッダー
        desc_header = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.3), Inches(1.0), Inches(0.3)
        )
        desc_header.text_frame.text = "説明"
        desc_header.text_frame.paragraphs[0].font.size = Pt(14)
        desc_header.text_frame.paragraphs[0].font.bold = True
        
        # 説明内容（左側の大きな枠）
        description = issue_data.get('description', '')
        if not description:
            description = "説明なし"
        
        # テキストが長い場合は適切に改行
        wrapped_description = self._wrap_text(description, 80)
        
        desc_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8), Inches(5.0), Inches(3.0)
        )
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = RGBColor(248, 249, 250)  # 薄いグレー背景
        desc_box.line.color.rgb = RGBColor(206, 212, 218)  # グレーボーダー
        
        desc_frame = desc_box.text_frame
        desc_frame.text = wrapped_description
        desc_frame.paragraphs[0].font.size = Pt(10)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(33, 37, 41)
        desc_frame.margin_left = Inches(0.1)
        desc_frame.margin_top = Inches(0.1)
        desc_frame.margin_right = Inches(0.1)
        desc_frame.margin_bottom = Inches(0.1)
    
    def _add_details_section(self, slide, issue_data: Dict):
        # 詳細ヘッダー
        details_header = slide.shapes.add_textbox(
            Inches(6.0), Inches(2.3), Inches(1.0), Inches(0.3)
        )
        details_header.text_frame.text = "詳細"
        details_header.text_frame.paragraphs[0].font.size = Pt(14)
        details_header.text_frame.paragraphs[0].font.bold = True
        
        # 詳細情報テーブル形式
        details = [
            ("作成者", issue_data.get('author', {}).get('name', '')),
            ("担当者", issue_data.get('assigned_to', {}).get('name', '') if issue_data.get('assigned_to') else ''),
            ("優先度", issue_data.get('priority', {}).get('name', '')),
            ("進捗率", f"{issue_data.get('done_ratio', 0)}%"),
            ("開始日", issue_data.get('start_date', '')),
            ("期限日", issue_data.get('due_date', '')),
            ("実績工数", f"{issue_data.get('spent_hours', 0)} 時間")
        ]
        
        y_pos = 2.8
        for label, value in details:
            if value:  # 値が空でない場合のみ表示
                # ラベル
                label_box = slide.shapes.add_textbox(
                    Inches(6.0), Inches(y_pos), Inches(1.5), Inches(0.25)
                )
                label_box.text_frame.text = label
                label_box.text_frame.paragraphs[0].font.size = Pt(10)
                label_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(108, 117, 125)
                
                # 値
                value_box = slide.shapes.add_textbox(
                    Inches(7.8), Inches(y_pos), Inches(1.7), Inches(0.25)
                )
                value_box.text_frame.text = str(value)
                value_box.text_frame.paragraphs[0].font.size = Pt(10)
                value_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 37, 41)
                
                y_pos += 0.3
    
    def _add_comments_section(self, slide, issue_data: Dict):
        # コメントヘッダー
        comment_header = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2), Inches(1.0), Inches(0.3)
        )
        comment_header.text_frame.text = "コメント"
        comment_header.text_frame.paragraphs[0].font.size = Pt(14)
        comment_header.text_frame.paragraphs[0].font.bold = True
        
        # コメント内容枠
        comment_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.7), Inches(9.0), Inches(0.8)
        )
        comment_box.fill.solid()
        comment_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        comment_box.line.color.rgb = RGBColor(206, 212, 218)
        comment_box.line.dash_style = 7  # 点線
        
        # コメント取得と表示
        comments_text = self._get_comments_text(issue_data)
        comment_frame = comment_box.text_frame
        comment_frame.text = comments_text
        comment_frame.paragraphs[0].font.size = Pt(9)
        comment_frame.paragraphs[0].font.color.rgb = RGBColor(108, 117, 125)
        comment_frame.margin_left = Inches(0.1)
        comment_frame.margin_top = Inches(0.1)
    
    def _add_footer_section(self, slide, issue_data: Dict):
        # 更新履歴フッター
        created_on = issue_data.get('created_on', '')
        updated_on = issue_data.get('updated_on', '')
        
        if created_on:
            created_on = created_on[:19].replace('T', ' ')
        if updated_on:
            updated_on = updated_on[:19].replace('T', ' ')
        
        footer_text = f"更新履歴 - 作成: {created_on}  最終更新: {updated_on}"
        
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.8), Inches(9.0), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = footer_text
        footer_frame.paragraphs[0].font.size = Pt(9)
        footer_frame.paragraphs[0].font.color.rgb = RGBColor(108, 117, 125)
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _get_comments_text(self, issue_data: Dict) -> str:
        journals = issue_data.get('journals', [])
        if not journals:
            return "コメントなし"
        
        comments = []
        for journal in journals:
            notes = journal.get('notes', '').strip()
            if notes:
                user_name = journal.get('user', {}).get('name', '不明なユーザー')
                created_on = journal.get('created_on', '')
                if created_on:
                    created_on = created_on[:19].replace('T', ' ')
                
                comment_text = f"[{user_name} - {created_on}]\n{notes}"
                comments.append(comment_text)
        
        if not comments:
            return "コメントなし"
        
        # 最新のコメントを1つだけ表示（スペースの都合上）
        latest_comment = comments[-1] if comments else "コメントなし"
        if len(latest_comment) > 200:
            latest_comment = latest_comment[:200] + "..."
        
        return latest_comment
    
    def _wrap_text(self, text: str, width: int) -> str:
        if not text:
            return ""
        
        lines = text.split('\n')
        wrapped_lines = []
        
        for line in lines:
            if len(line) <= width:
                wrapped_lines.append(line)
            else:
                wrapped_lines.extend(textwrap.wrap(line, width=width))
        
        return '\n'.join(wrapped_lines)
    
    def _save_to_bytes(self) -> bytes:
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()